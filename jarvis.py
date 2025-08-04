import openai
import json
import os
from reportlab.pdfgen import canvas
from docx import Document
from pptx import Presentation
from PIL import Image
import requests

# === Load API Key ===
with open("config.json") as f:
    config = json.load(f)
openai.api_key = config["OPENAI_API_KEY"]

# === Speak Function ===
def speak(text):
    print(f"JARVIS 🤖: {text}")

# === Text Commands ===
def handle_input(cmd):
    if cmd.startswith("pdf:"):
        content = cmd[4:].strip()
        filename = "output.pdf"
        c = canvas.Canvas(filename)
        c.drawString(100, 750, content)
        c.save()
        speak(f"PDF saved as {filename}.")

    elif cmd.startswith("docx:"):
        content = cmd[5:].strip()
        doc = Document()
        doc.add_paragraph(content)
        filename = "output.docx"
        doc.save(filename)
        speak(f"DOCX saved as {filename}.")

    elif cmd.startswith("ppt:"):
        content = cmd[4:].strip()
        ppt = Presentation()
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        slide.shapes.title.text = "Jarvis Presentation"
        slide.placeholders[1].text = content
        filename = "output.pptx"
        ppt.save(filename)
        speak(f"PPT saved as {filename}.")

    elif cmd.startswith("image:"):
        prompt = cmd[6:].strip()
        response = openai.Image.create(
            prompt=prompt,
            n=1,
            size="512x512"
        )
        image_url = response['data'][0]['url']
        img_data = requests.get(image_url).content
        with open("output.png", "wb") as handler:
            handler.write(img_data)
        speak("Image saved as output.png.")

    elif "who created you" in cmd.lower():
        speak("I am built by Shaan — SHAANX STUDIOS.")
    
    else:
        try:
            response = openai.ChatCompletion.create(
                model="gpt-4",
                messages=[{"role": "user", "content": cmd}]
            )
            speak(response['choices'][0]['message']['content'].strip())
        except Exception as e:
            speak(f"Error: {str(e)}")

# === Run Loop ===
if __name__ == "__main__":
    speak("Jarvis 3.0 activated. Type your command:")
    while True:
        user_input = input("You >> ")
        if user_input.lower() in ["exit", "quit"]:
            speak("Shutting down. Stay legendary.")
            break
        handle_input(user_input)
